#!/usr/bin/env python3
"""Сборка презентации presentation.pptx из метрик и графиков.

Числа (overall, difficulty, multiturn, document, fail-types) читаются из
metrics_axes.json / failure_analysis.json — так слайды всегда в синхроне с замерами.
Картинки берутся из slides/ и charts/. Нарратив фикса (стабильные числа эскалации) — в коде.

    python build_pptx.py
"""
import json
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Emu, Inches, Pt

ROOT = Path(__file__).resolve().parent
BLUE = RGBColor(0x2A, 0x7A, 0xDE)
DARK = RGBColor(0x22, 0x22, 0x22)
GREY = RGBColor(0x66, 0x66, 0x66)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
SW, SH = prs.slide_width, prs.slide_height
BLANK = prs.slide_layouts[6]


def slide():
    return prs.slides.add_slide(BLANK)


def textbox(s, left, top, width, height):
    tb = s.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    return tf


def set_title(s, text, sub=None):
    tf = textbox(s, Inches(0.6), Inches(0.35), Inches(12.1), Inches(1.0))
    p = tf.paragraphs[0]
    r = p.add_run(); r.text = text
    r.font.size = Pt(30); r.font.bold = True; r.font.color.rgb = DARK
    # синяя полоска под заголовком
    line = s.shapes.add_textbox(Inches(0.62), Inches(1.15), Inches(2.4), Inches(0.06))
    line.fill.solid() if False else None
    if sub:
        st = textbox(s, Inches(0.62), Inches(1.12), Inches(12.0), Inches(0.5))
        rp = st.paragraphs[0]; rr = rp.add_run(); rr.text = sub
        rr.font.size = Pt(14); rr.font.italic = True; rr.font.color.rgb = GREY


def bullets(s, items, left=Inches(0.7), top=Inches(1.7), width=Inches(12.0),
            height=Inches(5.3), size=18):
    tf = textbox(s, left, top, width, height)
    first = True
    for it, lvl in items:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.level = lvl
        p.space_after = Pt(6)
        # маркер
        prefix = "" if lvl == 0 else "– "
        r = p.add_run(); r.text = ("• " if lvl == 0 else prefix) + it
        r.font.size = Pt(size - lvl * 2)
        r.font.color.rgb = DARK if lvl == 0 else GREY
        if lvl == 0 and it.endswith(":"):
            r.font.bold = True


def add_image_fit(s, path, left, top, max_w, max_h):
    """Вписать картинку в EMU-бокс (left, top, max_w, max_h), сохранив пропорции, по центру."""
    from PIL import Image
    iw, ih = Image.open(path).size
    aspect = iw / ih
    if aspect > max_w / max_h:           # упираемся в ширину
        w = int(max_w); h = int(max_w / aspect)
    else:                                # упираемся в высоту
        h = int(max_h); w = int(max_h * aspect)
    pic = s.shapes.add_picture(str(path), 0, 0, width=w, height=h)
    pic.left = int(left + (max_w - w) / 2)
    pic.top = int(top + (max_h - h) / 2)
    return pic


def pct(x):
    return f"{x*100:.0f}%"


def panel(s, left, top, width, height, title, lines, accent, fill, mono_idx=()):
    """Цветная панель с заголовком и строками. mono_idx — индексы строк моноширинным."""
    box = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    box.fill.solid(); box.fill.fore_color.rgb = fill
    box.line.color.rgb = accent; box.line.width = Pt(1.5)
    box.shadow.inherit = False
    tf = box.text_frame; tf.word_wrap = True
    tf.margin_left = Inches(0.18); tf.margin_right = Inches(0.18)
    tf.margin_top = Inches(0.12)
    p = tf.paragraphs[0]; r = p.add_run(); r.text = title
    r.font.size = Pt(16); r.font.bold = True; r.font.color.rgb = accent
    for i, ln in enumerate(lines):
        pp = tf.add_paragraph(); pp.space_after = Pt(4)
        rr = pp.add_run(); rr.text = ln
        rr.font.size = Pt(12.5)
        rr.font.color.rgb = DARK
        if i in mono_idx:
            rr.font.name = "Courier New"; rr.font.size = Pt(12)
    return box


# --- данные из замеров ---
mx = json.loads((ROOT / "metrics_axes.json").read_text(encoding="utf-8"))
fa = json.loads((ROOT / "failure_analysis.json").read_text(encoding="utf-8"))
ax = mx["axes"]
overall = mx["overall"]; n_clean = mx["n"]


def rate_of(axis, key):
    d = ax[axis].get(key)
    return pct(d["rate"]) if d else "—"


# ============ Слайд 1 — Титул ============
s = slide()
tf = textbox(s, Inches(1.0), Inches(2.5), Inches(11.3), Inches(2.0))
p = tf.paragraphs[0]; r = p.add_run()
r.text = "Агент поддержки кредитования МСБ"
r.font.size = Pt(40); r.font.bold = True; r.font.color.rgb = DARK
p2 = tf.add_paragraph(); r2 = p2.add_run()
r2.text = "LangGraph (ReAct) · GigaChat-2 · FAISS RAG"
r2.font.size = Pt(22); r2.font.color.rgb = BLUE
p3 = tf.add_paragraph(); r3 = p3.add_run()
r3.text = "Отвечает по нормативке, работает с данными клиента, эскалирует и отказывает — с проверяемыми сигналами"
r3.font.size = Pt(16); r3.font.italic = True; r3.font.color.rgb = GREY

authors = textbox(s, Inches(1.0), Inches(5.4), Inches(11.3), Inches(1.2))
ap = authors.paragraphs[0]; ar = ap.add_run()
ar.text = "Фида А. Д.   ·   Сахаутдинова А. И."
ar.font.size = Pt(20); ar.font.bold = True; ar.font.color.rgb = DARK
ap2 = authors.add_paragraph(); ar2 = ap2.add_run()
ar2.text = "при участии Claude Code (Opus 4.8) — соавтор"
ar2.font.size = Pt(14); ar2.font.italic = True; ar2.font.color.rgb = GREY

# ============ Слайд 2 — Архитектура ============
s = slide()
set_title(s, "Архитектура: поток одного хода")
add_image_fit(s, ROOT / "slides/architecture.png", Inches(0.4), Inches(1.7),
              Inches(12.5), Inches(5.2))

# ============ Слайд 3 — Ключевые решения ============
s = slide()
set_title(s, "Ключевые решения реализации")
bullets(s, [
    ("RAG с иерархией заголовков:", 0),
    ("чанк несёт путь родительских заголовков (условия помнят продукт) + обзорный чанк; хранит source и номер пункта", 1),
    ("Безопасность данных в коде, не в промпте:", 0),
    ("клиентские тулы замкнуты на client_id сессии и не принимают id от модели → чужие данные закрыты архитектурно", 1),
    ("Машинно-проверяемая эскалация:", 0),
    ("escalate(trigger, summary) с триггерами sales/negative/security → детерминированный сигнал для метрик", 1),
    ("Рефлексия + долговременная память:", 0),
    ("самооценка ответа всей цепочки (fail-open) и remember/recall по client_id", 1),
])

# ============ Слайд 4 — Признаки системы ============
s = slide()
set_title(s, "Из чего состоит агент")
bullets(s, [
    ("ReAct agent loop — вызывает tool → получает результат → выбирает следующее действие, "
     "и так до ответа; всё автономно, без подтверждений промежуточных шагов", 0),
    ("RAG по нормативке — любой факт (продукт, ставка, требование) берётся из поиска FAISS, "
     "а не из памяти модели; чанки хранят источник и номер пункта для ссылок", 0),
    ("Клиентские инструменты — профиль / кредиты / заявки / расчёт досрочного погашения "
     "из БД, замкнуты на client_id сессии (чужие данные закрыты в коде)", 0),
    ("Долговременная память — remember/recall: хранит цель обращения, продукт, сумму/срок, "
     "поднимает контекст прошлых ходов и сессий клиента", 0),
    ("Self-reflection — доп. слой валидации: агент сам оценивает черновик ответа на соответствие "
     "регламенту и переигрывает при провале (fail-open)", 0),
    ("Эскалация как действие — escalate(trigger=sales|negative|security): реальный hand-off "
     "на человека + детерминированный сигнал для метрик", 0),
    ("Роль и ограничения — system-промпт по регламенту: что делает, чего НЕ делает (не обещает "
     "одобрение, не раскрывает скоринг), деловой тон", 0),
], top=Inches(1.6), size=17)

# ============ Слайд 5 — Как мерили ============
s = slide()
set_title(s, "Как мерили")
bullets(s, [
    (f"180 кейсов qa.jsonl, два независимых сигнала:", 0),
    ("escalation — детерминированно: сработал ли нужный триггер; не сработал ли ложно", 1),
    ("judge — LLM-судья (GigaChat) сверяет ответ с expected_behavior по сути", 1),
    ("Режем по осям: категория, сложность (easy/medium/hard), одно/многоходовость, документ", 0),
    ("Цель — показать «было → стало» после фиксов, а не максимизировать число", 0),
])

# ============ Слайд 6 — Результаты по осям ============
s = slide()
set_title(s, "Результаты по осям",
          sub=f"judge pass-rate {pct(overall)} на {n_clean} кейсах")
add_image_fit(s, ROOT / "charts/by_difficulty.png", Inches(0.3), Inches(1.9),
              Inches(6.4), Inches(4.6))
add_image_fit(s, ROOT / "charts/by_multiturn.png", Inches(6.7), Inches(1.9),
              Inches(6.4), Inches(4.6))
cap = textbox(s, Inches(0.7), Inches(6.6), Inches(12.0), Inches(0.6))
rc = cap.paragraphs[0].add_run()
rc.text = (f"Сложность монотонна: easy {rate_of('difficulty','easy')} → "
           f"medium {rate_of('difficulty','medium')} → hard {rate_of('difficulty','hard')}. "
           f"Многоходовые труднее одноходовых.")
rc.font.size = Pt(14); rc.font.color.rgb = GREY

# ============ Слайд 7 — Находка ============
s = slide()
set_title(s, "Находка: что eval поймал, а судья пропустил")
bullets(s, [
    ("Симптом: judge ставил «pass», но детерминированный escalation-сигнал sales = 17%", 0),
    ("Причина: GigaChat-2 печатал вызов инструмента ТЕКСТОМ —", 0),
    ("trigger=\"sales\", escalate(...), remember(...) вместо структурного tool-call", 1),
    ("похоже на ответ, но реального действия (hand-off / запись в память) НЕ происходит", 1),
    ("Масштаб: 11 из 180 кейсов (6%), стохастично между прогонами", 0),
    ("Мораль: одного LLM-судьи мало — нужен детерминированный сигнал по факту вызова инструмента", 0),
])

# ============ Слайд 8 — Что пофиксили (было/стало) ============
s = slide()
set_title(s, "Что именно пофиксили",
          sub="Клиент пишет: «Хочу оформить кредит на развитие» → агент должен передать заявку человеку (escalate sales)")
RED = RGBColor(0xC0, 0x39, 0x2B)
GREEN = RGBColor(0x2E, 0x86, 0x4F)
panel(s, Inches(0.6), Inches(1.95), Inches(6.0), Inches(3.9),
      "БЫЛО  ❌",
      ['В промпте триггер описан как ярлык:',
       '   1. trigger="sales" — клиент хочет…',
       'Модель копирует ярлык в ТЕКСТ ответа:',
       '   trigger="sales"  Понял ваше намерение…',
       'Инструмент escalate НЕ вызван:',
       '   • эскалации в системе нет → оператор',
       '     не подключён, заявка повисла',
       '   • клиент видит «сырой» trigger="sales"',
       'Метрика escalation_sales = 17%'],
      RED, RGBColor(0xFC, 0xF3, 0xF2), mono_idx=(1, 3, 5))
panel(s, Inches(6.9), Inches(1.95), Inches(6.0), Inches(3.9),
      "СТАЛО  ✅",
      ['Фикс 1 — промпт: триггер как ДЕЙСТВИЕ +',
       '   «эскалация только вызовом инструмента,',
       '    не пиши trigger= текстом»',
       'Фикс 2 — guard в agent.py (детерминир.):',
       '   видит текстовый вызов в ответе →',
       '   • поднимает реальную escalate(sales)',
       '   • прячет «сырой» вызов, шлёт hand-off:',
       '     «Передаю обращение специалисту…»',
       'escalation_sales 17% → 39% → 44%'],
      GREEN, RGBColor(0xF1, 0xFA, 0xF4), mono_idx=(2,))
note = textbox(s, Inches(0.6), Inches(6.05), Inches(12.3), Inches(1.2))
for i, ln in enumerate([
    "Guard не зависит от «послушности» модели и ТОЛЬКО добавляет эскалации (никогда не убирает).",
    "Бонус-фикс: GIGACHAT_CREDENTIALS — авто-обновление токена, длинные прогоны eval не падают на «Token expired».",
]):
    p = note.paragraphs[0] if i == 0 else note.add_paragraph()
    r = p.add_run(); r.text = "• " + ln; r.font.size = Pt(13); r.font.color.rgb = GREY

# ============ Слайд 9 — Было → стало ============
s = slide()
set_title(s, "Было → стало: эскалация (детерминир. сигнал)")
add_image_fit(s, ROOT / "charts/escalation_beforeafter.png", Inches(0.5), Inches(1.7),
              Inches(7.8), Inches(5.2))
tf = textbox(s, Inches(8.6), Inches(2.0), Inches(4.4), Inches(4.8))
for it, lvl in [
    ("sales: 17% → 39% → 44%", 0),
    ("negative: 88% → 56% → 69%", 0),
    ("(baseline → +промпт → +guard)", 1),
    ("Guard даёт монотонный плюс", 0),
    ("Честно: negative baseline 88% —", 0),
    ("tool-calling GigaChat стохастичен между прогонами (88%↔56%); guard ограничивает просадку, но не перебивает дисперсию", 1),
]:
    p = tf.paragraphs[0] if not tf.paragraphs[0].runs else tf.add_paragraph()
    p.space_after = Pt(8)
    rr = p.add_run(); rr.text = ("• " if lvl == 0 else "  – ") + it
    rr.font.size = Pt(16 if lvl == 0 else 13)
    rr.font.color.rgb = DARK if lvl == 0 else GREY

# ============ Слайды 10-12 — Примеры трейсов (чат-интерфейс) ============
import math

TRACES = {t["id"]: t for t in json.loads((ROOT / "traces.json").read_text(encoding="utf-8"))}
GREEN2 = RGBColor(0x2E, 0x86, 0x4F)
ORANGE = RGBColor(0xCF, 0x8B, 0x1F)


def _clip(t, n):
    t = str(t).replace("\n", " ").strip()
    return t if len(t) <= n else t[:n] + "…"


def bubble(s, top, kind, header, body):
    fill, accent, mono, width_in, right = {
        "user":   (RGBColor(0xDC, 0xEB, 0xFB), BLUE, False, 6.4, True),
        "tool":   (RGBColor(0xFF, 0xF1, 0xD6), ORANGE, True, 8.8, False),
        "result": (RGBColor(0xED, 0xF1, 0xF3), GREY, True, 8.8, False),
        "reason": (RGBColor(0xF2, 0xF2, 0xF2), GREY, False, 8.8, False),
        "final":  (RGBColor(0xE6, 0xF7, 0xEC), GREEN2, False, 9.4, False),
    }[kind]
    width = Inches(width_in)
    left = (SW - width - Inches(0.5)) if right else Inches(0.5)
    cpl = int(width_in * (7.0 if mono else 7.8))
    lines = sum(max(1, math.ceil(len(p) / max(1, cpl))) for p in str(body).split("\n"))
    height = Inches(0.38 + lines * 0.205)
    box = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, int(left), int(top),
                             int(width), int(height))
    box.fill.solid(); box.fill.fore_color.rgb = fill
    box.line.color.rgb = accent; box.line.width = Pt(1.0); box.shadow.inherit = False
    tf = box.text_frame; tf.word_wrap = True
    tf.margin_left = Inches(0.12); tf.margin_right = Inches(0.12)
    tf.margin_top = Inches(0.05); tf.margin_bottom = Inches(0.05)
    p = tf.paragraphs[0]; r = p.add_run(); r.text = header
    r.font.size = Pt(11); r.font.bold = True; r.font.color.rgb = accent
    pp = tf.add_paragraph(); rr = pp.add_run(); rr.text = str(body)
    rr.font.size = Pt(11 if mono else 12.5); rr.font.color.rgb = DARK
    if mono:
        rr.font.name = "Courier New"
    return int(top) + int(height) + Inches(0.09)


def chat_slide(tid, title):
    t = TRACES[tid]
    s = slide()
    sub = f"{t['category']} · канал {t.get('channel') or '—'}"
    if t.get("client_id"):
        sub += f" · client_id {t['client_id']}"
    set_title(s, title, sub=sub)
    top = int(Inches(1.5))
    limit = int(Inches(7.15))
    top = bubble(s, top, "user", "👤 Клиент", _clip(t["question"], 240))
    steps = t["trace"]
    last_ai = max((i for i, st in enumerate(steps)
                   if st["type"] == "AIMessage" and (st.get("content") or "").strip()),
                  default=-1)
    for i, st in enumerate(steps):
        if top > limit:
            bubble(s, top, "reason", "…", "(трейс обрезан для слайда)")
            break
        typ = st["type"]
        if typ in ("SystemMessage", "HumanMessage"):
            continue
        if typ == "AIMessage":
            for tc in st.get("tool_calls") or []:
                args = json.dumps(tc["args"], ensure_ascii=False)
                top = bubble(s, top, "tool", "🔧 вызов инструмента",
                             f"{tc['name']}({_clip(args, 100)})")
            content = (st.get("content") or "").strip()
            if content:
                if i == last_ai:
                    top = bubble(s, top, "final", "🤖 Ответ агента", _clip(content, 300))
                else:
                    top = bubble(s, top, "reason", "🤖 рассуждает", _clip(content, 150))
        elif typ == "ToolMessage":
            top = bubble(s, top, "result", f"⚙️ результат · {st.get('tool_name', 'tool')}",
                         _clip(st.get("content", ""), 160))
    if t.get("escalations"):
        trig = ", ".join(e["trigger"] for e in t["escalations"])
        cy = min(top, int(Inches(7.05)))
        chip = textbox(s, Inches(0.5), cy, Inches(8.5), Inches(0.35))
        rc = chip.paragraphs[0].add_run()
        rc.text = f"✅ эскалация зафиксирована: trigger = {trig}"
        rc.font.size = Pt(12); rc.font.bold = True; rc.font.color.rgb = GREEN2


chat_slide("Q-001", "Пример 1 · RAG по нормативке")
chat_slide("Q-061", "Пример 2 · Клиентские тулы + расчёт досрочного погашения")
chat_slide("Q-091", "Пример 3 · Эскалация (sales) → реальный hand-off")

# ============ Слайд 13 — Чего не хватает / что улучшить ============
s = slide()
set_title(s, "Чего не хватает и что улучшить")
bullets(s, [
    ("Память — уровень 1 (локальный JSON):", 0),
    ("до уровня 2 нужен внешний долговременный стор («цифровой двойник» клиента) с поиском по смыслу", 1),
    ("Корень проблемы — стохастичность tool-calling GigaChat:", 0),
    ("guard смягчает, но надёжнее — structured-output / function-calling режим + ретраи невалидных вызовов", 1),
    ("Краевые случаи (edge_*) и transactional — слабее всего:", 0),
    ("отдельный eligibility-шаг (проверка сегмента/требований ДО рекомендации) + few-shot примеры отказов; точный график платежей для расчёта ДП", 1),
    ("Оценка качества — один прогон судьи:", 0),
    ("judge даёт дисперсию; нужно ≥3 прогонов с усреднением, межсудейское согласие, и метрики precision/recall эскалации отдельно", 1),
    ("Надёжность и стоимость:", 0),
    ("фолбэк на локальные эмбеддинги при недоступности API; рефлексия удваивает LLM-вызовы — включать выборочно", 1),
], top=Inches(1.55), size=17)

out = ROOT / "presentation.pptx"
prs.save(str(out))
print(f"Готово: {out}  ({len(prs.slides._sldIdLst)} слайдов)")
